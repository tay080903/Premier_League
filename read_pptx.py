import collections 
import collections.abc
import pptx
from pptx import Presentation

# load
prs = Presentation('EPL Nations Hub Project Proposal.pptx')

# text
text_runs = []
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_runs.append(shape.text)

with open('proposal.txt', 'w', encoding='utf-8') as f:
    f.write("\n".join(text_runs))
